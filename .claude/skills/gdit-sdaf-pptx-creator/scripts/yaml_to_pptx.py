#!/usr/bin/env python3
# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx", "pyyaml"]
# ///
"""YAML-to-PPTX converter. Reads a YAML content file and produces a .pptx presentation.

Usage:
    python3 yaml_to_pptx.py <input.yaml> [--output <file.pptx>]
"""

import argparse
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BUILTIN_COLORS = {
    "white": "#FFFFFF",
    "black": "#000000",
    "green": "#27AE60",
    "red": "#C0392B",
}


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Theme:
    def __init__(self, data: dict):
        d = data or {}
        self.primary = d.get("primary", "#1B2A4A")
        self.secondary = d.get("secondary", "#2E5C9A")
        self.accent = d.get("accent", "#E86C00")
        self.background = d.get("background", "#FFFFFF")
        self.text = d.get("text", "#333333")
        self.muted = d.get("muted", "#666666")
        self.light_bg = d.get("light_bg", "#F0F2F5")

    def resolve(self, name_or_hex: str | None) -> str:
        if not name_or_hex:
            return self.accent
        if name_or_hex.startswith("#"):
            return name_or_hex
        mapping = {
            "primary": self.primary,
            "secondary": self.secondary,
            "accent": self.accent,
            "background": self.background,
            "text": self.text,
            "muted": self.muted,
            "light_bg": self.light_bg,
            "green": BUILTIN_COLORS["green"],
            "red": BUILTIN_COLORS["red"],
            "white": BUILTIN_COLORS["white"],
            "black": BUILTIN_COLORS["black"],
        }
        return mapping.get(name_or_hex, self.accent)


# ── Drawing helpers ──

def add_bg(slide, color_hex: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def add_rect(slide, left, top, width, height, color_hex):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color_hex)
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color="#FFFFFF", bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = hex_to_rgb(color)
    p.font.bold = bold
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, size=16, color="#333333", spacing=Pt(6)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.font.size = Pt(size)
        p.font.color.rgb = hex_to_rgb(color)
        p.space_after = spacing
    return tf


def add_header_bar(slide, title, theme):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), theme.primary)
    add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6),
             title, size=30, color="#FFFFFF", bold=True)


def add_footer(slide, text, theme):
    if text:
        add_text(slide, Inches(0.5), SLIDE_H - Inches(0.45), Inches(6), Inches(0.3),
                 text, size=9, color=theme.muted)


def add_metric_box(slide, left, top, width, height, value, label, val_color, theme):
    add_rect(slide, left, top, width, height, "#FFFFFF")
    add_text(slide, left, top + Inches(0.15), width, Inches(0.6),
             value, size=32, color=val_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.75), width, Inches(0.5),
             label, size=13, color=theme.muted, align=PP_ALIGN.CENTER)


# ── Slide renderers ──

def render_title(sd, prs, theme, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme.primary)
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.15), SLIDE_W, Inches(0.15), theme.accent)
    add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
             sd.get("title", ""), size=44, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    if sd.get("subtitle"):
        add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 sd["subtitle"], size=24, color=theme.secondary, align=PP_ALIGN.CENTER)
    if sd.get("tagline"):
        add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
                 sd["tagline"], size=16, color=theme.muted, align=PP_ALIGN.CENTER)
    if sd.get("bottom_text"):
        add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 sd["bottom_text"], size=14, color=theme.muted, align=PP_ALIGN.CENTER)


def render_content(sd, prs, theme, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme.background)
    add_header_bar(slide, sd.get("title", ""), theme)
    if sd.get("intro"):
        add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
                 sd["intro"], size=17, color=theme.text)
    body_top = Inches(2.0) if sd.get("intro") else Inches(1.4)
    if sd.get("body"):
        add_bullets(slide, Inches(0.8), body_top, Inches(11.5), Inches(4.5),
                    sd["body"], size=15, color=theme.text)
    add_footer(slide, footer, theme)


def render_two_column(sd, prs, theme, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme.background)
    add_header_bar(slide, sd.get("title", ""), theme)
    if sd.get("intro"):
        add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
                 sd["intro"], size=17, color=theme.text)
    col_top = Inches(2.0) if sd.get("intro") else Inches(1.4)
    col_w, col_h = Inches(5.8), Inches(4.8) if not sd.get("intro") else Inches(4.2)

    for i, side in enumerate(["left", "right"]):
        col = sd.get(side, {})
        x = Inches(0.6) if i == 0 else Inches(6.9)
        style = col.get("style", "light")
        bg = theme.light_bg if style == "light" else theme.primary
        txt = theme.text if style == "light" else "#FFFFFF"
        head_c = theme.primary if style == "light" else theme.accent

        add_rect(slide, x, col_top, col_w, col_h, bg)
        if col.get("heading"):
            add_text(slide, x + Inches(0.3), col_top + Inches(0.15), Inches(5.2), Inches(0.4),
                     col["heading"], size=18, color=head_c, bold=True)
        if col.get("items"):
            add_bullets(slide, x + Inches(0.3), col_top + Inches(0.65), Inches(5.2), col_h - Inches(1.0),
                        col["items"], size=14, color=txt, spacing=Pt(4))
        if col.get("footer_text"):
            add_text(slide, x + Inches(0.3), col_top + col_h - Inches(0.7), Inches(5.2), Inches(0.6),
                     col["footer_text"], size=14, color=theme.muted if style == "light" else theme.accent)
    add_footer(slide, footer, theme)


def render_metrics(sd, prs, theme, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme.background)
    add_header_bar(slide, sd.get("title", ""), theme)

    metrics = sd.get("metrics", [])
    count = len(metrics)
    if count > 0:
        box_w = min(Inches(2.3), Inches(12) / count - Inches(0.2))
        box_h = Inches(1.2)
        total_w = count * box_w + (count - 1) * Inches(0.3)
        start_x = (SLIDE_W - total_w) / 2
        for i, m in enumerate(metrics):
            vc = theme.resolve(m.get("color"))
            add_metric_box(slide, start_x + i * (box_w + Inches(0.3)), Inches(1.4),
                           box_w, box_h, str(m.get("value", "")), str(m.get("label", "")), vc, theme)

    if sd.get("heading"):
        add_text(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.5),
                 sd["heading"], size=20, color=theme.primary, bold=True)
    body_top = Inches(3.5) if sd.get("heading") else Inches(3.0)
    if sd.get("body"):
        add_bullets(slide, Inches(0.8), body_top, Inches(11.5), Inches(3.5),
                    sd["body"], size=15, color=theme.text, spacing=Pt(5))
    add_footer(slide, footer, theme)


def render_grid(sd, prs, theme, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme.background)
    add_header_bar(slide, sd.get("title", ""), theme)

    cells = sd.get("cells", [])
    cols = sd.get("columns", 3)
    if not cells:
        add_footer(slide, footer, theme)
        return

    rows = (len(cells) + cols - 1) // cols
    cw = (Inches(12.2) - (cols - 1) * Inches(0.35)) / cols
    ch = min(Inches(2.2), (Inches(5.0) - (rows - 1) * Inches(0.3)) / rows)

    if sd.get("intro"):
        add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
                 sd["intro"], size=16, color=theme.text)

    grid_top = Inches(1.7) if sd.get("intro") else Inches(1.4)
    for idx, cell in enumerate(cells):
        r, c = divmod(idx, cols)
        x = Inches(0.6) + c * (cw + Inches(0.35))
        y = grid_top + r * (ch + Inches(0.3))
        add_rect(slide, x, y, cw, ch, theme.light_bg)
        if cell.get("heading"):
            add_text(slide, x + Inches(0.2), y + Inches(0.1), cw - Inches(0.4), Inches(0.35),
                     cell["heading"], size=15, color=theme.primary, bold=True)
        if cell.get("items"):
            add_bullets(slide, x + Inches(0.2), y + Inches(0.5), cw - Inches(0.4), ch - Inches(0.6),
                        cell["items"], size=11, color=theme.text, spacing=Pt(2))

    if sd.get("bottom_text"):
        add_text(slide, Inches(0.8), grid_top + rows * (ch + Inches(0.3)), Inches(11), Inches(0.5),
                 sd["bottom_text"], size=14, color=theme.muted)
    add_footer(slide, footer, theme)


def render_call_to_action(sd, prs, theme, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme.primary)
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.15), SLIDE_W, Inches(0.15), theme.accent)

    add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.7),
             sd.get("title", ""), size=36, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)

    columns = sd.get("columns", [])
    if columns:
        count = len(columns)
        col_w = min(Inches(3.6), (Inches(11.5) - (count - 1) * Inches(0.4)) / count)
        total_w = count * col_w + (count - 1) * Inches(0.4)
        start_x = (SLIDE_W - total_w) / 2
        for i, col in enumerate(columns):
            x = start_x + i * (col_w + Inches(0.4))
            add_rect(slide, x, Inches(1.8), col_w, Inches(3.5), theme.secondary)
            if col.get("heading"):
                add_text(slide, x, Inches(2.0), col_w, Inches(0.4),
                         col["heading"], size=20, color=theme.accent, bold=True, align=PP_ALIGN.CENTER)
            if col.get("items"):
                add_bullets(slide, x + Inches(0.3), Inches(2.5), col_w - Inches(0.6), Inches(2.5),
                            col["items"], size=14, color="#FFFFFF", spacing=Pt(6))

    if sd.get("closing"):
        add_text(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.6),
                 sd["closing"], size=18, color="#FFFFFF", align=PP_ALIGN.CENTER)
    if sd.get("bottom_text"):
        add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 sd["bottom_text"], size=16, color=theme.muted, align=PP_ALIGN.CENTER)


def parse_drawio_labels(filepath: str) -> tuple[str, list[str]]:
    """Extract diagram name and text labels from a .drawio XML file."""
    import re
    from pathlib import Path
    text = Path(filepath).read_text()
    # Get diagram name from <diagram> element
    name_match = re.search(r'<diagram[^>]*name="([^"]*)"', text)
    name = name_match.group(1) if name_match else Path(filepath).stem
    # Extract value attributes (text labels) from mxCell elements
    values = re.findall(r'value="([^"]+)"', text)
    # Clean HTML entities and tags, deduplicate, skip empty
    labels = []
    seen = set()
    for v in values:
        clean = re.sub(r'<[^>]+>', ' ', v).replace('&#xa;', ' ').replace('&amp;', '&')
        clean = re.sub(r'\s+', ' ', clean).strip()
        if clean and clean not in seen and len(clean) > 1:
            seen.add(clean)
            labels.append(clean)
    return name, labels


def render_diagram(sd, prs, theme, footer):
    """Render an appendix slide for a diagram. Embeds PNG if available, falls back to text labels."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme.background)
    title = sd.get("title", "Diagram")
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), theme.secondary)
    add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6),
             f"Appendix: {title}", size=28, color="#FFFFFF", bold=True)

    filepath = sd.get("file", "")
    fpath = Path(filepath) if filepath else None

    # Check for image: explicit image field, or .png/.svg next to .drawio
    image_path = None
    if sd.get("image"):
        img = Path(sd["image"])
        if img.exists():
            image_path = img
    elif fpath and fpath.suffix == ".drawio":
        for ext in (".png", ".svg", ".jpg"):
            candidate = fpath.with_suffix(ext)
            if candidate.exists():
                image_path = candidate
                break

    # Auto-convert .drawio to PNG if no image found
    if not image_path and fpath and fpath.exists() and fpath.suffix == ".drawio":
        try:
            from importlib.util import spec_from_file_location, module_from_spec
            converter_path = Path(__file__).parent / "drawio_to_png.py"
            if converter_path.exists():
                spec = spec_from_file_location("drawio_to_png", converter_path)
                mod = module_from_spec(spec)
                spec.loader.exec_module(mod)
                import tempfile
                tmp_png = Path(tempfile.mktemp(suffix=".png"))
                if mod.drawio_to_png(filepath, str(tmp_png)):
                    image_path = tmp_png
        except Exception:
            pass  # Fall through to text fallback

    if image_path:
        # Embed the image centered on the slide
        add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
                 f"Source: {filepath}", size=10, color=theme.muted)
        from PIL import Image as PILImage
        with PILImage.open(str(image_path)) as img:
            iw, ih = img.size  # pixels
        # Convert to inches at 150 DPI (matches ImageMagick density)
        iw_in = iw / 150.0
        ih_in = ih / 150.0
        max_w_in, max_h_in = 11.5, 5.2
        scale = min(max_w_in / iw_in, max_h_in / ih_in, 1.0)
        w = Inches(iw_in * scale)
        h = Inches(ih_in * scale)
        left = (SLIDE_W - w) // 2
        slide.shapes.add_picture(str(image_path), left, Inches(1.6), w, h)
    elif fpath and fpath.exists() and fpath.suffix == ".drawio":
        # Final fallback: extract text labels from .drawio XML
        name, labels = parse_drawio_labels(filepath)
        add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
                 f"Diagram: {name}", size=16, color=theme.primary, bold=True)
        add_text(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.3),
                 f"Source: {filepath}  (install ImageMagick for rendered images)", size=11, color=theme.muted)
        if labels:
            mid = (len(labels) + 1) // 2 if len(labels) > 12 else len(labels)
            add_bullets(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5),
                        labels[:mid], size=12, color=theme.text, spacing=Pt(2))
            if len(labels) > 12:
                add_bullets(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(4.5),
                            labels[mid:], size=12, color=theme.text, spacing=Pt(2))
    elif filepath:
        add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
                 f"Source: {filepath}", size=14, color=theme.muted)
        if sd.get("description"):
            add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
                     sd["description"], size=16, color=theme.text)
    add_footer(slide, footer, theme)


# ── Dispatcher and CLI ──

RENDERERS = {
    "title": render_title,
    "content": render_content,
    "two-column": render_two_column,
    "metrics": render_metrics,
    "grid": render_grid,
    "call-to-action": render_call_to_action,
    "diagram": render_diagram,
}


def load_yaml(path: str) -> dict:
    with open(path) as f:
        return yaml.safe_load(f)


def create_presentation(data: dict) -> tuple[Presentation, int]:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    pdata = data.get("presentation", {})
    theme = Theme(pdata.get("theme"))
    footer = pdata.get("footer", "")
    slides = list(data.get("slides", []))

    # Auto-expand appendix.diagrams into diagram slides
    appendix = data.get("appendix", {})
    for diag in appendix.get("diagrams", []):
        slides.append({
            "type": "diagram",
            "title": diag.get("title", Path(diag.get("file", "")).stem),
            "file": diag.get("file", ""),
            "image": diag.get("image", ""),
            "description": diag.get("description", ""),
        })

    for i, sd in enumerate(slides):
        stype = sd.get("type", "content")
        renderer = RENDERERS.get(stype)
        if renderer is None:
            print(f"ERROR: Unknown slide type '{stype}' at slide index {i}", file=sys.stderr)
            sys.exit(1)
        renderer(sd, prs, theme, footer)
        # Add speaker notes if present
        notes_text = sd.get("notes")
        if notes_text:
            slide = prs.slides[len(prs.slides) - 1]
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    return prs, len(slides)


def main():
    parser = argparse.ArgumentParser(description="Convert YAML to PPTX presentation")
    parser.add_argument("input", help="YAML content file")
    parser.add_argument("--output", "-o", help="Output .pptx path (default: derived from title)")
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"ERROR: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    data = load_yaml(str(input_path))
    prs, count = create_presentation(data)

    if args.output:
        out = Path(args.output)
    else:
        title = data.get("presentation", {}).get("title", "presentation")
        safe = title.replace(" ", "-").replace("/", "-")[:60]
        out = input_path.parent / f"{safe}.pptx"

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"✓ Created: {out}")
    print(f"  Slides: {count}")
    print(f"  Size: {out.stat().st_size / 1024:.0f} KB")


if __name__ == "__main__":
    main()
